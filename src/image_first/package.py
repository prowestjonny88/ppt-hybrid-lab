import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from src.ir.runtime import canonical_hash, load_json

SLIDES = [
    ("problem-hook", "problem_hook.jpg", "problem_hook.v1.json"),
    ("how-it-works", "how_it_works.jpg", "how_it_works.v1.json"),
    ("validation-traction", "validation_traction.jpg", "validation_traction.v1.json"),
]
TARGET_RATIO = 16 / 9
MAX_RELATIVE_RATIO_ERROR = 0.03


def _name_shape(shape, name):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("shape has no cNvPr")


def _validate_image(path):
    with Image.open(path) as image:
        width, height = image.size
        image_format = image.format
    if image_format not in {"JPEG", "JPG"}:
        raise RuntimeError(f"image-first Gemini asset must be JPEG: {path} is {image_format}")
    ratio = width / height
    relative_error = abs(ratio - TARGET_RATIO) / TARGET_RATIO
    if relative_error > MAX_RELATIVE_RATIO_ERROR:
        raise RuntimeError(
            f"image-first asset is too far from requested 16:9: {path} is {width}x{height} "
            f"(ratio={ratio:.6f}, relative_error={relative_error:.4%})"
        )
    return width, height, ratio


def _center_crop_picture_to_ratio(picture, source_ratio, target_ratio):
    """Crop source pixels symmetrically while retaining the original JPEG bytes in PPTX."""
    if abs(source_ratio - target_ratio) < 1e-9:
        return {"left": 0.0, "right": 0.0, "top": 0.0, "bottom": 0.0}
    if source_ratio > target_ratio:
        # Source is slightly too wide. Keep all height and crop equal horizontal fractions.
        keep_width_fraction = target_ratio / source_ratio
        crop = (1.0 - keep_width_fraction) / 2.0
        picture.crop_left = crop
        picture.crop_right = crop
        return {"left": crop, "right": crop, "top": 0.0, "bottom": 0.0}
    # Source is slightly too tall. Keep all width and crop equal vertical fractions.
    keep_height_fraction = source_ratio / target_ratio
    crop = (1.0 - keep_height_fraction) / 2.0
    picture.crop_top = crop
    picture.crop_bottom = crop
    return {"left": 0.0, "right": 0.0, "top": crop, "bottom": crop}


def build_image_first(root, assets_dir, output_pptx, manifest_path):
    root = Path(root)
    assets_dir = Path(assets_dir)
    deck = load_json(root / "experiment/queuezero/deck_system.stage3.json")
    prs = Presentation()
    prs.slide_width = Inches(deck["slide_size"]["width_in"])
    prs.slide_height = Inches(deck["slide_size"]["height_in"])
    blank = prs.slide_layouts[6]
    manifest = {
        "schema_version": "stage3-image-first-package-v2",
        "variant": "image_first",
        "placement_policy": "center_crop_source_without_resampling",
        "slides": [],
    }
    for slide_id, filename, semantics_name in SLIDES:
        image_path = assets_dir / filename
        if not image_path.exists():
            raise RuntimeError(f"missing generated image-first asset: {image_path}")
        width, height, ratio = _validate_image(image_path)
        semantics = load_json(root / "experiment/queuezero/slide_semantics" / semantics_name)
        slide = prs.slides.add_slide(blank)
        picture = slide.shapes.add_picture(str(image_path), 0, 0, prs.slide_width, prs.slide_height)
        crop = _center_crop_picture_to_ratio(picture, ratio, TARGET_RATIO)
        _name_shape(picture, f"oxq:{slide_id}:full_slide:image_first")
        manifest["slides"].append({
            "slide_id": slide_id,
            "image": str(image_path),
            "pixel_size": [width, height],
            "source_ratio": ratio,
            "target_ratio": TARGET_RATIO,
            "crop_fractions": crop,
            "semantic_hash": canonical_hash(semantics),
            "ppt_shape_name": picture.name,
            "editability_class": "full_slide_raster_baseline",
        })
    output_pptx = Path(output_pptx)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    manifest_path = Path(manifest_path)
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(json.dumps(manifest, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    return output_pptx, manifest
