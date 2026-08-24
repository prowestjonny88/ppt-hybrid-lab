#!/usr/bin/env python3
"""Render Stage 4 compiled layout solutions into editable PowerPoint objects."""

from __future__ import annotations

import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.ir.runtime import load_json, object_index, resolved_object_text
from src.visual_ir.style_runtime import resolve_style_profile

SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
STYLE_PATH = "experiment/queuezero/style_profiles/queuezero_hackathon_v0.json"


def _rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _name_shape(shape, name: str):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError(f"unable to name shape {name}: cNvPr not found")


def _to_emu_box(prs: Presentation, box):
    x, y, w, h = box
    return (
        int(x * prs.slide_width), int(y * prs.slide_height),
        int(w * prs.slide_width), int(h * prs.slide_height),
    )


def _shape_signature(shape):
    text = shape.text if getattr(shape, "has_text_frame", False) else ""
    payload = json.dumps({
        "shape_id": shape.shape_id,
        "name": shape.name,
        "text": text,
        "bounds": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
    }, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _text_for(placement, semantic_objects, semantics):
    obj = semantic_objects[placement["semantic_object_id"]]
    source = placement["text_source"]
    if source == "content":
        return resolved_object_text(obj, semantics)
    if source == "label":
        return obj.get("label", "")
    raise RuntimeError(f"unsupported text_source {source!r}")


def _add_text(slide, prs, placement, style, semantic_objects, semantics):
    x, y, w, h = _to_emu_box(prs, placement["box"])
    shape = slide.shapes.add_textbox(x, y, w, h)
    _name_shape(shape, f"oxq:{semantics['slide_id']}:{placement['semantic_object_id']}:{placement['part']}")
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[placement["valign"]]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[placement["align"]]
    p.space_before = p.space_after = Pt(0)
    run = p.add_run(); run.text = _text_for(placement, semantic_objects, semantics)
    run.font.name = style["typography"]["primary_family_preference"][0]
    run.font.size = Pt(placement["font_size_pt"])
    run.font.bold = placement["font_weight"] >= 600
    run.font.color.rgb = _rgb(style["palette"][placement["color_token"]])
    return shape


def _add_decor(slide, prs, decor, style, slide_id):
    x, y, w, h = _to_emu_box(prs, decor["box"])
    kind = decor["kind"]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE if kind == "rect" else MSO_SHAPE.OVAL, x, y, w, h) if kind in {"rect", "ellipse"} else None
    if shape is None:
        raise RuntimeError(f"unsupported decoration kind {kind!r}")
    _name_shape(shape, f"oxq:{slide_id}:decor:{decor['decor_id']}")
    shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(style["palette"][decor["fill_token"]])
    if decor.get("line_token"):
        shape.line.color.rgb = _rgb(style["palette"][decor["line_token"]]); shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _add_line(slide, prs, placement, style, semantics):
    x, y, w, h = _to_emu_box(prs, placement["box"])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, max(h, Pt(1)))
    _name_shape(shape, f"oxq:{semantics['slide_id']}:{placement['semantic_object_id']}:{placement['part']}")
    shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(style["palette"][placement["color_token"]])
    shape.line.fill.background()
    return shape


def _add_picture(slide, prs, root, placement, style, semantics):
    x, y, w, h = _to_emu_box(prs, placement["box"])
    asset = root / placement["asset_path"]
    if asset.is_file():
        shape = slide.shapes.add_picture(str(asset), x, y, w, h)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(style["palette"][placement["fallback_fill_token"]])
        shape.line.fill.background()
    _name_shape(shape, f"oxq:{semantics['slide_id']}:{placement['semantic_object_id']}:{placement['part']}")
    return shape, asset.is_file()


def render_layout_slide(root: Path, prs: Presentation, layout_solution: dict):
    root = Path(root)
    if layout_solution.get("schema_version") != "stage4-layout-solution-v0":
        raise RuntimeError("unsupported layout solution schema")
    semantics = load_json(root / layout_solution["semantic_file"])
    style, profile, language = resolve_style_profile(root, root / STYLE_PATH)
    if layout_solution.get("style_profile_id") != profile["profile_id"]:
        raise RuntimeError("layout solution/style profile mismatch")
    if layout_solution.get("design_language_id") != language["design_language_id"]:
        raise RuntimeError("layout solution/design language mismatch")
    semantic_objects = object_index(semantics)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    canvas = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _name_shape(canvas, f"oxq:{semantics['slide_id']}:decor:canvas")
    canvas.fill.solid(); canvas.fill.fore_color.rgb = _rgb(style["palette"][layout_solution["canvas"]["background_token"]]); canvas.line.fill.background()
    realization_objects = []
    decoration_objects = [{"name": canvas.name, "shape_id": str(canvas.shape_id), "content_hash": _shape_signature(canvas)}]

    for decor in layout_solution["decorations"]:
        shape = _add_decor(slide, prs, decor, style, semantics["slide_id"])
        decoration_objects.append({"decor_id": decor["decor_id"], "name": shape.name, "shape_id": str(shape.shape_id), "content_hash": _shape_signature(shape)})

    for placement in layout_solution["placements"]:
        kind = placement["kind"]
        asset_present = None
        if kind == "text":
            shape = _add_text(slide, prs, placement, style, semantic_objects, semantics)
            lane, fidelity = "native", "semantic_and_editable"
        elif kind == "line":
            shape = _add_line(slide, prs, placement, style, semantics)
            lane, fidelity = "native", "semantic_and_editable"
        elif kind == "picture":
            shape, asset_present = _add_picture(slide, prs, root, placement, style, semantics)
            lane, fidelity = "hybrid_asset", "replaceable_asset"
        else:
            raise RuntimeError(f"unsupported placement kind {kind!r}")
        record = {
            "semantic_object_id": placement["semantic_object_id"], "part": placement["part"],
            "kind": kind, "ppt_shape_name": shape.name, "ppt_object_id": str(shape.shape_id),
            "bounds_emu": {"x": int(shape.left), "y": int(shape.top), "w": int(shape.width), "h": int(shape.height)},
            "content_hash": _shape_signature(shape), "render_lane": lane, "fidelity": fidelity,
        }
        if asset_present is not None:
            record["asset_present"] = asset_present; record["asset_path"] = placement["asset_path"]
        realization_objects.append(record)

    return {
        "slide_id": semantics["slide_id"],
        "semantic_hash": layout_solution["semantic_hash"],
        "visual_ir_hash": layout_solution["visual_ir_hash"],
        "design_language_id": layout_solution["design_language_id"],
        "design_language_hash": layout_solution["design_language_hash"],
        "style_profile_id": layout_solution["style_profile_id"],
        "style_profile_hash": layout_solution["style_profile_hash"],
        "resolved_style_hash": layout_solution["resolved_style_hash"],
        "archetype_id": layout_solution["archetype_id"],
        "variant": layout_solution["variant"],
        "routing_trace": layout_solution.get("routing_trace"),
        "solver": layout_solution.get("solver"),
        "compiler": layout_solution["compiler"],
        "objects": realization_objects,
        "decorations": decoration_objects,
    }


def new_presentation():
    prs = Presentation(); prs.slide_width = Inches(SLIDE_W_IN); prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def render_compiled_sample(root: Path, layout_solution: dict, output_pptx: Path, realization_path: Path):
    root = Path(root); output_pptx = Path(output_pptx); realization_path = Path(realization_path)
    if layout_solution.get("slide_id") != "validation-traction":
        raise RuntimeError("representative sample renderer requires validation-traction")
    prs = new_presentation(); realization = render_layout_slide(root, prs, layout_solution)
    output_pptx.parent.mkdir(parents=True, exist_ok=True); prs.save(output_pptx)
    realization = {"schema_version": "stage4-realization-v0", "pptx": str(output_pptx), **realization}
    realization_path.parent.mkdir(parents=True, exist_ok=True)
    realization_path.write_text(json.dumps(realization, indent=2) + "\n", encoding="utf-8")
    return output_pptx, realization
