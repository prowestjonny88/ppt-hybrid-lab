import hashlib
import json
import os
import subprocess
import sys
from html import escape
from pathlib import Path

from pptx import Presentation

from src.ir.runtime import load_json, object_index

PPT_MASTER_COMMIT = "65bb2eca59a36270819caba377097910c4466c6e"
SLIDE_W = 1280
SLIDE_H = 720


def _shape_name(shape, name):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("shape has no cNvPr")


def _token(deck, name):
    return deck["tokens"][name]


def _svg_text(x, y, text, *, element_id, size=16, fill="#111827", weight=500):
    return (
        f'<text id="{escape(element_id)}" x="{x}" y="{y}" '
        f'font-family="Arial" font-size="{size}" font-weight="{weight}" '
        f'fill="{fill}">{escape(text)}</text>'
    )


def write_how_it_works_svg(root, svg_path):
    root = Path(root)
    semantics = load_json(root / "experiment/queuezero/slide_semantics/how_it_works.v1.json")
    deck = load_json(root / "experiment/queuezero/deck_system.stage3.json")
    objects = object_index(semantics)
    region = next(r for r in semantics["regions"] if r["region_id"] == "diagram_region")["rect"]

    rx, ry = region["x"] * SLIDE_W, region["y"] * SLIDE_H
    rw, rh = region["w"] * SLIDE_W, region["h"] * SLIDE_H
    nodes = {
        "node_camera": (0.01, 0.29, 0.20, 0.30),
        "node_queue_estimator": (0.27, 0.29, 0.20, 0.30),
        "node_wait_predictor": (0.53, 0.29, 0.20, 0.30),
        "node_decision": (0.79, 0.24, 0.20, 0.40),
    }
    accent = _token(deck, "accent.primary")
    accent_on = _token(deck, "accent.on")
    surface = _token(deck, "surface.raised")
    text_primary = _token(deck, "text.primary")
    line = _token(deck, "line.subtle")

    body = []
    for object_id, (nx, ny, nw, nh) in nodes.items():
        x, y, w, h = rx + nx * rw, ry + ny * rh, nw * rw, nh * rh
        is_accent = object_id == "node_wait_predictor"
        body.append(
            f'<rect id="{object_id}__shape" x="{x:.2f}" y="{y:.2f}" '
            f'width="{w:.2f}" height="{h:.2f}" rx="12" ry="12" '
            f'fill="{accent if is_accent else surface}" '
            f'stroke="{accent if is_accent else line}" stroke-width="1.5"/>'
        )
        label = objects[object_id]["content"]
        if object_id == "node_decision":
            parts = ["Go now / choose", "another cafeteria"]
            for idx, part in enumerate(parts):
                body.append(_svg_text(
                    x + w * 0.10,
                    y + h * (0.43 + idx * 0.20),
                    part,
                    element_id=f"{object_id}__label{idx+1}",
                    size=13,
                    fill=accent_on if is_accent else text_primary,
                    weight=600 if is_accent else 500,
                ))
        else:
            body.append(_svg_text(
                x + w * 0.12,
                y + h * 0.57,
                label,
                element_id=f"{object_id}__label",
                size=14,
                fill=accent_on if is_accent else text_primary,
                weight=600 if is_accent else 500,
            ))

    connectors = [
        ("connector_camera_queue", 0.21, 0.44, 0.27, 0.44),
        ("connector_queue_prediction", 0.47, 0.44, 0.53, 0.44),
        ("connector_prediction_decision", 0.73, 0.44, 0.79, 0.44),
    ]
    for object_id, x1n, y1n, x2n, y2n in connectors:
        x1, y1 = rx + x1n * rw, ry + y1n * rh
        x2, y2 = rx + x2n * rw, ry + y2n * rh
        body.append(
            f'<line id="{object_id}__line" x1="{x1:.2f}" y1="{y1:.2f}" '
            f'x2="{x2 - 8:.2f}" y2="{y2:.2f}" stroke="{accent}" stroke-width="2.5"/>'
        )
        body.append(
            f'<polygon id="{object_id}__arrow" points="{x2-9:.2f},{y2-6:.2f} '
            f'{x2:.2f},{y2:.2f} {x2-9:.2f},{y2+6:.2f}" fill="{accent}"/>'
        )

    svg = (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{SLIDE_W}" height="{SLIDE_H}" '
        f'viewBox="0 0 {SLIDE_W} {SLIDE_H}">\n'
        + "\n".join(body)
        + "\n</svg>\n"
    )
    svg_path = Path(svg_path)
    svg_path.parent.mkdir(parents=True, exist_ok=True)
    svg_path.write_text(svg, encoding="utf-8")
    return svg_path


def _run(cmd, cwd):
    cwd = Path(cwd)
    env = os.environ.copy()
    scripts_root = cwd / "skills/ppt-master/scripts"
    existing = env.get("PYTHONPATH", "")
    env["PYTHONPATH"] = str(scripts_root) + (os.pathsep + existing if existing else "")
    result = subprocess.run(cmd, cwd=str(cwd), text=True, capture_output=True, env=env)
    if result.returncode != 0:
        raise RuntimeError(
            "PPT Master command failed:\n"
            + " ".join(map(str, cmd))
            + "\n--- stdout ---\n" + result.stdout[-8000:]
            + "\n--- stderr ---\n" + result.stderr[-8000:]
        )
    return result


def _trace_events(value):
    if isinstance(value, dict):
        if isinstance(value.get("events"), list):
            yield from value["events"]
        for key, child in value.items():
            if key != "events":
                yield from _trace_events(child)
    elif isinstance(value, list):
        for child in value:
            yield from _trace_events(child)


def _semantic_part(svg_id):
    if "__" not in svg_id:
        return svg_id, "main"
    object_id, part = svg_id.split("__", 1)
    return object_id, part


def stamp_trace_identity(pptx_path, trace_path, output_path=None):
    pptx_path = Path(pptx_path)
    output_path = Path(output_path or pptx_path)
    trace = json.loads(Path(trace_path).read_text(encoding="utf-8"))
    prs = Presentation(pptx_path)
    if len(prs.slides) != 1:
        raise RuntimeError(f"expected exactly one PPT Master slide, got {len(prs.slides)}")
    slide = prs.slides[0]
    shapes_by_id = {str(shape.shape_id): shape for shape in slide.shapes}
    mapped = {}
    for event in _trace_events(trace):
        if event.get("decision") != "native" or not event.get("id") or event.get("shape_id") is None:
            continue
        svg_id = str(event["id"])
        shape_id = str(event["shape_id"])
        if shape_id not in shapes_by_id:
            continue
        object_id, part = _semantic_part(svg_id)
        shape = shapes_by_id[shape_id]
        name = f"oxq:how-it-works:{object_id}:{part}"
        _shape_name(shape, name)
        mapped.setdefault(object_id, []).append(shape_id)
    if not mapped:
        raise RuntimeError("PPT Master trace produced no native semantic mappings")
    prs.save(output_path)
    return mapped


def compile_how_it_works(root, ppt_master_root, workspace):
    root = Path(root)
    ppt_master_root = Path(ppt_master_root)
    workspace = Path(workspace)
    svg_path = write_how_it_works_svg(root, workspace / "svg_output/queuezero_hiw.svg")
    validation = workspace / "validation"
    validation.mkdir(parents=True, exist_ok=True)

    checker = ppt_master_root / "skills/ppt-master/scripts/svg_quality_checker.py"
    cli = ppt_master_root / "skills/ppt-master/scripts/svg_to_pptx/pptx_package/cli.py"
    if not checker.exists() or not cli.exists():
        raise RuntimeError("pinned PPT Master checkout is incomplete")

    actual_commit = _run(["git", "rev-parse", "HEAD"], ppt_master_root).stdout.strip()
    if actual_commit != PPT_MASTER_COMMIT:
        raise RuntimeError(f"PPT Master commit mismatch: expected {PPT_MASTER_COMMIT}, got {actual_commit}")

    _run([
        sys.executable, str(checker), str(workspace),
        "--quick-generate", "--stage", "final", "--json",
    ], ppt_master_root)

    out_pptx = workspace / "out.pptx"
    trace_path = workspace / "trace.json"
    _run([
        sys.executable, str(cli), str(workspace),
        "--quick-generate", "--no-animations",
        "-o", str(out_pptx), "--conversion-trace", str(trace_path), "-q",
    ], ppt_master_root)

    report = validation / "out.report.json"
    if not out_pptx.exists() or not trace_path.exists() or not report.exists():
        raise RuntimeError("PPT Master did not produce required adapter artifacts")
    stamped = workspace / "out.stamped.pptx"
    mapped = stamp_trace_identity(out_pptx, trace_path, stamped)
    return {
        "svg": svg_path,
        "pptx": stamped,
        "trace": trace_path,
        "postflight": report,
        "mapped": mapped,
        "svg_sha256": hashlib.sha256(svg_path.read_bytes()).hexdigest(),
        "ppt_master_commit": actual_commit,
    }
