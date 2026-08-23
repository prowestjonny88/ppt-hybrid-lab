import hashlib
import json
from pathlib import Path

from PIL import Image, ImageDraw
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.ir.runtime import object_index, region_index


def _rgb(value):
    value = value.lstrip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _name_shape(shape, name):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("shape has no cNvPr")


def ensure_mock_ui(path):
    path = Path(path)
    if path.exists():
        return path
    path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (720, 1280), "#F7F8FA")
    d = ImageDraw.Draw(image)
    d.rounded_rectangle((35, 35, 685, 1245), radius=38, fill="#FFFFFF", outline="#D1D5DB", width=4)
    d.rounded_rectangle((75, 95, 645, 250), radius=20, fill="#E8F0FF")
    d.text((105, 155), "QueueZero — live cafeteria waits", fill="#111827")
    rows = [(320, "North cafeteria", "8 min"), (525, "Central cafeteria", "18 min"), (730, "South cafeteria", "11 min")]
    for y, label, wait in rows:
        d.rounded_rectangle((85, y, 635, y + 165), radius=20, fill="#FFFFFF", outline="#E5E7EB", width=3)
        d.text((120, y + 55), label, fill="#111827")
        d.text((510, y + 55), wait, fill="#2563EB")
    d.rounded_rectangle((155, 1000, 565, 1120), radius=30, fill="#2563EB")
    d.text((270, 1045), "Best option", fill="#FFFFFF")
    image.save(path)
    return path


class Stage3Emitter:
    def __init__(self, semantics, deck_system, slide, variant="native_vector"):
        self.semantics = semantics
        self.deck = deck_system
        self.slide = slide
        self.variant = variant
        self.objects = object_index(semantics)
        self.regions = region_index(semantics)
        self.slide_w = int(Inches(deck_system["slide_size"]["width_in"]))
        self.slide_h = int(Inches(deck_system["slide_size"]["height_in"]))
        self._parts = {}

    def color(self, token):
        return _rgb(self.deck["tokens"][token])

    def region(self, region_id):
        r = self.regions[region_id]["rect"]
        return int(r["x"]*self.slide_w), int(r["y"]*self.slide_h), int(r["w"]*self.slide_w), int(r["h"]*self.slide_h)

    def box(self, region_id, x=0, y=0, w=1, h=1):
        rx, ry, rw, rh = self.region(region_id)
        return int(rx+x*rw), int(ry+y*rh), int(w*rw), int(h*rh)

    def _record(self, object_id, shape, part, fidelity="semantic_and_editable", tokens=None, lane="native"):
        name = f"oxq:{self.semantics['slide_id']}:{object_id}:{part}"
        _name_shape(shape, name)
        text = shape.text if getattr(shape, "has_text_frame", False) else ""
        signature = json.dumps({
            "name": name,
            "shape_id": shape.shape_id,
            "text": text,
            "bounds": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
        }, sort_keys=True, separators=(",", ":"))
        entry = {
            "ppt_object_id": str(shape.shape_id),
            "ppt_shape_name": name,
            "bounds_emu": {"x": int(shape.left), "y": int(shape.top), "w": int(shape.width), "h": int(shape.height)},
            "content_hash": hashlib.sha256(signature.encode()).hexdigest(),
        }
        rec = self._parts.setdefault(object_id, {
            "semantic_object_id": object_id,
            "render_lane": lane,
            "parts": [],
            "token_refs": list(tokens or self.objects.get(object_id, {}).get("token_refs", [])),
            "fidelity": fidelity,
            "fit_result": "fit",
            "fallback": None,
            "realization_rev": 1,
        })
        rec["parts"].append(entry)
        return shape

    def text(self, object_id, text, region_id, box=(0,0,1,1), role=None, part="text", align="left", valign="top", token=None, size=None, bold=None):
        x,y,w,h = self.box(region_id, *box)
        shape = self.slide.shapes.add_textbox(x,y,w,h)
        tf = shape.text_frame
        tf.clear(); tf.word_wrap = True
        tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
        tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run(); run.text = text
        role = role or self.objects.get(object_id, {}).get("role", "body_text")
        defaults = self.deck["role_defaults"].get(role, self.deck["role_defaults"]["body_text"])
        type_spec = self.deck["type_scale"][defaults["type_token"]]
        run.font.size = Pt(size or type_spec["size_pt"])
        run.font.bold = bool(type_spec["weight"] >= 600) if bold is None else bold
        token = token or (self.objects.get(object_id, {}).get("token_refs") or defaults.get("token_refs") or ["text.primary"])[0]
        run.font.color.rgb = self.color(token)
        return self._record(object_id, shape, part, tokens=[token])

    def rect(self, object_id, region_id, box, fill="surface.raised", line="line.subtle", part="shape", rounded=True, lane="native"):
        x,y,w,h = self.box(region_id, *box)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x,y,w,h)
        shape.fill.solid(); shape.fill.fore_color.rgb = self.color(fill)
        if line:
            shape.line.color.rgb = self.color(line); shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return self._record(object_id, shape, part, tokens=[fill] + ([line] if line else []), lane=lane)

    def metric_card(self, object_id, value, label, region_id, box):
        x,y,w,h = box
        self.rect(object_id, region_id, box, part="card")
        self.text(object_id, value, region_id, (x+0.06*w,y+0.10*h,0.88*w,0.48*h), role="metric", part="value", align="center", valign="middle", token="accent.primary")
        self.text(object_id, label, region_id, (x+0.06*w,y+0.62*h,0.88*w,0.22*h), role="metric_label", part="label", align="center", valign="middle", token="text.secondary")

    def node(self, object_id, text, region_id, box, accent=False, lane="native"):
        fill = "accent.primary" if accent else "surface.raised"
        self.rect(object_id, region_id, box, fill=fill, line=None if accent else "line.subtle", part="node", lane=lane)
        x,y,w,h = box
        token = "accent.on" if accent else "text.primary"
        shape = self.text(object_id, text, region_id, (x+0.06*w,y+0.20*h,0.88*w,0.58*h), role="diagram_node", part="label", align="center", valign="middle", token=token, size=12, bold=accent)
        if lane != "native":
            # The temporary structured benchmark emitter can visually realize the lane,
            # but the canonical SVG adapter will replace this before hybrid scoring.
            self._parts[object_id]["render_lane"] = lane
        return shape

    def connector(self, object_id, region_id, x1,y1,x2,y2, lane="native"):
        rx,ry,rw,rh = self.region(region_id)
        ax,ay,bx,by = int(rx+x1*rw),int(ry+y1*rh),int(rx+x2*rw),int(ry+y2*rh)
        shape = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax,ay,bx,by)
        shape.line.color.rgb = self.color("accent.primary"); shape.line.width = Pt(2)
        self._record(object_id, shape, "line", tokens=["accent.primary"], lane=lane)
        tw,th = int(0.016*rw),int(0.05*rh)
        tri = self.slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, bx-tw, by-th//2, tw, th)
        tri.fill.solid(); tri.fill.fore_color.rgb = self.color("accent.primary"); tri.line.fill.background()
        self._record(object_id, tri, "arrow", tokens=["accent.primary"], lane=lane)

    def picture(self, object_id, source, region_id, lane="native"):
        x,y,w,h = self.box(region_id)
        source = ensure_mock_ui(source)
        shape = self.slide.shapes.add_picture(str(source), x,y,w,h)
        return self._record(object_id, shape, "picture", fidelity="editable_picture_slot", lane=lane)

    def queue_scene(self, object_id, region_id):
        self.rect(object_id, region_id, (0.04,0.06,0.92,0.86), fill="surface.raised", part="scene", lane="native")
        self.rect(object_id, region_id, (0.62,0.18,0.25,0.18), fill="accent.primary", line=None, part="counter", lane="native")
        self.text(object_id, "CAFETERIA", region_id, (0.64,0.215,0.21,0.08), role="source_note", part="counter_label", align="center", valign="middle", token="accent.on", size=8, bold=True)
        for idx,(px,py) in enumerate([(0.18,0.64),(0.32,0.57),(0.46,0.51),(0.59,0.45)],1):
            x,y,w,h = self.box(region_id,px,py,0.075,0.13)
            s = self.slide.shapes.add_shape(MSO_SHAPE.OVAL,x,y,w,h)
            s.fill.solid(); s.fill.fore_color.rgb = self.color("accent.primary" if idx==1 else "text.secondary"); s.line.fill.background()
            self._record(object_id,s,f"person{idx}",tokens=["accent.primary" if idx==1 else "text.secondary"])
        self.text(object_id,"unpredictable peak queue",region_id,(0.13,0.80,0.68,0.08),role="source_note",part="caption",align="center",token="text.secondary",size=9)

    def realization(self):
        records=[]
        for object_id,rec in self._parts.items():
            aggregate="|".join(part["content_hash"] for part in rec["parts"])
            rec["pptx_content_hash"] = hashlib.sha256(aggregate.encode()).hexdigest()
            records.append(rec)
        return {"schema_version":"stage3-realization-v1","slide_id":self.semantics["slide_id"],"variant":self.variant,"objects":records}
