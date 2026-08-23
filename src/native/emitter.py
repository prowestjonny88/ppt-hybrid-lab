import hashlib
import json
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.ir.runtime import object_index, region_index


def rgb(hex_value):
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_shape_name(shape, name):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("cNvPr not found")


def ensure_mock_product_ui(path):
    path = Path(path)
    if path.exists():
        return path
    path.parent.mkdir(parents=True, exist_ok=True)
    canvas = Image.new("RGB", (720, 1280), "#F7F8FA")
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle((42, 45, 678, 1235), radius=40, fill="#FFFFFF", outline="#D1D5DB", width=4)
    draw.rounded_rectangle((80, 110, 640, 260), radius=22, fill="#E8F0FF")
    draw.text((115, 155), "QueueZero", fill="#111827")
    cards = [
        (90, 330, 630, 500, "North Cafeteria", "8 min"),
        (90, 535, 630, 705, "Central Cafeteria", "18 min"),
        (90, 740, 630, 910, "South Cafeteria", "11 min"),
    ]
    for x1, y1, x2, y2, label, wait in cards:
        draw.rounded_rectangle((x1, y1, x2, y2), radius=20, fill="#FFFFFF", outline="#E5E7EB", width=3)
        draw.text((125, y1 + 48), label, fill="#111827")
        draw.text((500, y1 + 48), wait, fill="#2563EB")
    draw.rounded_rectangle((145, 1015, 575, 1125), radius=28, fill="#2563EB")
    draw.text((260, 1054), "Best option", fill="#FFFFFF")
    canvas.save(path)
    return path


class NativeEmitter:
    def __init__(self, semantics, deck_system, slide, variant):
        self.semantics = semantics
        self.deck_system = deck_system
        self.slide = slide
        self.variant = variant
        self.objects = object_index(semantics)
        self.regions = region_index(semantics)
        self.records = {}
        self.slide_w = int(slide.part.package.presentation_part.presentation.slide_width)
        self.slide_h = int(slide.part.package.presentation_part.presentation.slide_height)

    def token(self, name):
        return self.deck_system["tokens"][name]

    def _region_rect_emu(self, region_id):
        rect = self.regions[region_id]["rect"]
        return (
            int(rect["x"] * self.slide_w),
            int(rect["y"] * self.slide_h),
            int(rect["w"] * self.slide_w),
            int(rect["h"] * self.slide_h),
        )

    def box(self, region_id, x=0.0, y=0.0, w=1.0, h=1.0):
        rx, ry, rw, rh = self._region_rect_emu(region_id)
        return (
            int(rx + x * rw),
            int(ry + y * rh),
            int(w * rw),
            int(h * rh),
        )

    def _shape_text(self, shape):
        return shape.text if getattr(shape, "has_text_frame", False) else ""

    def _register(self, object_id, shape, part="main", fidelity="semantic_and_editable", token_refs=None, fit_result=None):
        name = f"oxq:{self.semantics['slide_id']}:{object_id}:{part}"
        set_shape_name(shape, name)
        payload = {
            "name": name,
            "shape_id": str(shape.shape_id),
            "text": self._shape_text(shape),
            "bounds": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
        }
        content_hash = hashlib.sha256(json.dumps(payload, sort_keys=True).encode()).hexdigest()
        rec = self.records.setdefault(object_id, {
            "semantic_object_id": object_id,
            "render_lane": "native",
            "ppt_object_ids": [],
            "ppt_shape_names": [],
            "bounds_emu": [],
            "token_refs": list(token_refs or self.objects.get(object_id, {}).get("token_refs", [])),
            "fidelity": fidelity,
            "fit_result": fit_result,
            "fallback": None,
            "shape_content_hashes": [],
            "realization_rev": 1,
        })
        rec["ppt_object_ids"].append(str(shape.shape_id))
        rec["ppt_shape_names"].append(name)
        rec["bounds_emu"].append({"x": int(shape.left), "y": int(shape.top), "w": int(shape.width), "h": int(shape.height)})
        rec["shape_content_hashes"].append(content_hash)
        return shape

    def add_text(self, object_id, text, region_id, box=(0, 0, 1, 1), role=None, align="left", valign="top", part="text", font_size=None, bold=None, token_ref=None):
        x, y, w, h = self.box(region_id, *box)
        shape = self.slide.shapes.add_textbox(x, y, w, h)
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        role = role or self.objects.get(object_id, {}).get("role", "body_text")
        defaults = self.deck_system["role_defaults"].get(role, self.deck_system["role_defaults"].get("body_text"))
        type_token = defaults["type_token"]
        type_spec = self.deck_system["type_scale"][type_token]
        run.font.size = Pt(font_size or type_spec["size_pt"])
        run.font.bold = bool(type_spec["weight"] >= 600) if bold is None else bold
        color_token = token_ref or (self.objects.get(object_id, {}).get("token_refs") or defaults.get("token_refs") or ["text.primary"])[0]
        run.font.color.rgb = rgb(self.token(color_token))
        return self._register(object_id, shape, part=part, token_refs=[color_token], fit_result="fit")

    def add_rect(self, object_id, region_id, box=(0, 0, 1, 1), fill_token="surface.raised", radius=True, part="shape", line_token=None):
        x, y, w, h = self.box(region_id, *box)
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = self.slide.shapes.add_shape(shape_type, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(self.token(fill_token))
        if line_token:
            shape.line.color.rgb = rgb(self.token(line_token))
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return self._register(object_id, shape, part=part, token_refs=[fill_token] + ([line_token] if line_token else []))

    def add_metric_card(self, object_id, value, label, region_id, box):
        x, y, w, h = box
        self.add_rect(object_id, region_id, (x, y, w, h), fill_token="surface.raised", radius=True, part="card", line_token="line.subtle")
        self.add_text(object_id, value, region_id, (x + 0.07*w, y + 0.13*h, 0.86*w, 0.46*h), role="metric", align="center", valign="middle", part="value", token_ref="accent.primary")
        self.add_text(object_id, label, region_id, (x + 0.05*w, y + 0.62*h, 0.90*w, 0.22*h), role="metric_label", align="center", valign="middle", part="label", token_ref="text.secondary")

    def add_connector(self, object_id, region_id, x1, y1, x2, y2, token_ref="accent.primary"):
        rx, ry, rw, rh = self._region_rect_emu(region_id)
        ax, ay, bx, by = int(rx + x1*rw), int(ry + y1*rh), int(rx + x2*rw), int(ry + y2*rh)
        line = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax, ay, bx, by)
        line.line.color.rgb = rgb(self.token(token_ref))
        line.line.width = Pt(2)
        self._register(object_id, line, part="line", token_refs=[token_ref])
        # Small native arrowhead triangle, kept under the same semantic connector identity.
        tri_w, tri_h = int(0.018*rw), int(0.055*rh)
        tri = self.slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, bx - tri_w, by - tri_h//2, tri_w, tri_h)
        tri.fill.solid(); tri.fill.fore_color.rgb = rgb(self.token(token_ref)); tri.line.fill.background()
        self._register(object_id, tri, part="arrow", token_refs=[token_ref])

    def add_picture(self, object_id, source, region_id, box=(0, 0, 1, 1), part="picture"):
        x, y, w, h = self.box(region_id, *box)
        source = ensure_mock_product_ui(source)
        pic = self.slide.shapes.add_picture(str(source), x, y, w, h)
        return self._register(object_id, pic, part=part, fidelity="editable_picture_slot", fit_result="fit")

    def add_queue_scene(self, object_id, region_id):
        # Native-vector visual baseline: a counter, service point, queue lane and people.
        self.add_rect(object_id, region_id, (0.06, 0.10, 0.88, 0.78), fill_token="surface.raised", radius=True, part="scene_bg", line_token="line.subtle")
        self.add_rect(object_id, region_id, (0.62, 0.18, 0.24, 0.18), fill_token="accent.primary", radius=True, part="counter")
        self.add_text(object_id, "CAFETERIA", region_id, (0.64, 0.21, 0.20, 0.10), role="caption", align="center", valign="middle", part="counter_label", token_ref="accent.on", font_size=9, bold=True)
        positions = [(0.20,0.62),(0.34,0.56),(0.48,0.50),(0.61,0.45)]
        for i, (px, py) in enumerate(positions, 1):
            x, y, w, h = self.box(region_id, px, py, 0.08, 0.13)
            person = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
            person.fill.solid(); person.fill.fore_color.rgb = rgb(self.token("accent.primary") if i == 1 else self.token("text.secondary")); person.line.fill.background()
            self._register(object_id, person, part=f"person{i}", token_refs=["accent.primary" if i == 1 else "text.secondary"])
        self.add_text(object_id, "unpredictable peak queue", region_id, (0.13, 0.78, 0.65, 0.08), role="caption", align="center", part="scene_caption", token_ref="text.secondary", font_size=9)

    def realization(self):
        for rec in self.records.values():
            combined = "|".join(rec.pop("shape_content_hashes"))
            rec["pptx_content_hash"] = hashlib.sha256(combined.encode()).hexdigest()
        return {
            "schema_version": "stage3-realization-v1",
            "slide_id": self.semantics["slide_id"],
            "variant": self.variant,
            "objects": list(self.records.values()),
        }
