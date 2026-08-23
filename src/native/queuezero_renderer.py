import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from src.ir.runtime import load_json, object_index, resolve_template, resolved_object_text, canonical_hash
from src.native.stage3_emitter import Stage3Emitter

SLIDES = [
    ("problem_hook.v1.json", "problem_hook.native_vector.json"),
    ("how_it_works.v1.json", "how_it_works.native_vector.json"),
    ("validation_traction.v1.json", "validation_traction.native_vector.json"),
]


def _text(semantics, object_id):
    return resolved_object_text(object_index(semantics)[object_id], semantics)


def _label(semantics, object_id):
    return object_index(semantics)[object_id].get("label", "")


def _render_problem(em):
    s = em.semantics
    em.text("title", _text(s,"title"), "title_band", (0,0,1,0.52), role="title", part="title")
    em.text("subtitle", _text(s,"subtitle"), "title_band", (0,0.58,1,0.32), role="subtitle", part="subtitle")
    em.metric_card("metric_wait", _text(s,"metric_wait"), _label(s,"metric_wait"), "problem_copy", (0.00,0.04,0.65,0.36))
    em.text("pain_annotation", _text(s,"pain_annotation"), "problem_copy", (0.00,0.50,0.92,0.38), role="annotation", part="annotation")
    em.queue_scene("hero_visual_slot", "hero_frame")
    em.text("source_note", _text(s,"source_note"), "footer_region", role="source_note", part="source")


def _render_how(em, root):
    s = em.semantics
    em.text("title", _text(s,"title"), "title_band", (0,0,1,0.52), role="title", part="title")
    em.text("subtitle", _text(s,"subtitle"), "title_band", (0,0.58,1,0.32), role="subtitle", part="subtitle")
    node_boxes = {
        "node_camera": (0.01,0.29,0.20,0.30),
        "node_queue_estimator": (0.27,0.29,0.20,0.30),
        "node_wait_predictor": (0.53,0.29,0.20,0.30),
        "node_decision": (0.79,0.24,0.20,0.40),
    }
    for object_id, box in node_boxes.items():
        em.node(object_id, _text(s,object_id), "diagram_region", box, accent=(object_id=="node_wait_predictor"))
    em.connector("connector_camera_queue", "diagram_region", 0.21,0.44,0.27,0.44)
    em.connector("connector_queue_prediction", "diagram_region", 0.47,0.44,0.53,0.44)
    em.connector("connector_prediction_decision", "diagram_region", 0.73,0.44,0.79,0.44)
    asset = next(a for a in s["assets"] if a["semantic_object_id"] == "screenshot_main")
    source = root / asset["current_instance"]["source"]
    em.picture("screenshot_main", source, "screenshot_region")
    em.text("source_note", _text(s,"source_note"), "footer_region", role="source_note", part="source")


def _render_validation(em):
    s = em.semantics
    em.text("title", _text(s,"title"), "title_band", (0,0,1,0.50), role="title", part="title")
    em.text("subtitle", _text(s,"subtitle"), "title_band", (0,0.56,1,0.34), role="subtitle", part="subtitle")
    boxes = {
        "metric_mae": (0.02,0.04,0.45,0.39),
        "metric_weekly_intent": (0.53,0.04,0.45,0.39),
        "metric_students": (0.02,0.53,0.45,0.31),
        "metric_cafeterias": (0.53,0.53,0.45,0.31),
    }
    for object_id, box in boxes.items():
        em.metric_card(object_id, _text(s,object_id), _label(s,object_id), "proof_region", box)
    em.rect("pilot_gate", "pilot_region", (0,0,1,1), fill="surface.raised", line="status.pending", part="gate")
    em.text("pilot_gate", _text(s,"pilot_gate"), "pilot_region", (0.04,0.15,0.92,0.70), role="annotation", part="label", align="center", valign="middle", token="status.pending", bold=True)
    em.text("source_note", _text(s,"source_note"), "footer_region", role="source_note", part="source")


def build_native_vector(root, output_pptx, realization_dir):
    root = Path(root)
    deck_system = load_json(root / "experiment/queuezero/deck_system.stage3.json")
    prs = Presentation()
    prs.slide_width = Inches(deck_system["slide_size"]["width_in"])
    prs.slide_height = Inches(deck_system["slide_size"]["height_in"])
    blank = prs.slide_layouts[6]
    realizations = []

    for semantic_name, plan_name in SLIDES:
        semantics = load_json(root / "experiment/queuezero/slide_semantics" / semantic_name)
        plan = load_json(root / "experiment/queuezero/render_plans" / plan_name)
        if plan["semantic_file"] != f"experiment/queuezero/slide_semantics/{semantic_name}":
            raise RuntimeError(f"render plan semantic mismatch for {semantic_name}")
        slide = prs.slides.add_slide(blank)
        em = Stage3Emitter(semantics, deck_system, slide, variant="native_vector")
        if semantics["slide_id"] == "problem-hook":
            _render_problem(em)
        elif semantics["slide_id"] == "how-it-works":
            _render_how(em, root)
        elif semantics["slide_id"] == "validation-traction":
            _render_validation(em)
        else:
            raise RuntimeError(f"unknown slide {semantics['slide_id']}")
        realization = em.realization()
        realization["semantic_hash"] = canonical_hash(semantics)
        realization["render_plan"] = f"experiment/queuezero/render_plans/{plan_name}"
        realizations.append(realization)

    output_pptx = Path(output_pptx)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)

    realization_dir = Path(realization_dir)
    realization_dir.mkdir(parents=True, exist_ok=True)
    for realization in realizations:
        path = realization_dir / f"{realization['slide_id']}.native_vector.json"
        path.write_text(json.dumps(realization, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    return output_pptx, realizations
