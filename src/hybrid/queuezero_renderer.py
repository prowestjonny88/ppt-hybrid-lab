import hashlib
import json
import tempfile
from pathlib import Path

from pptx import Presentation

from src.ir.runtime import canonical_hash, load_json, object_index, resolved_object_text
from src.native.queuezero_renderer import _label, _render_validation, _text
from src.native.stage3_emitter import Stage3Emitter
from src.svg_lane.ppt_master_adapter import compile_how_it_works


def _shape_signature(shape):
    text = shape.text if getattr(shape, "has_text_frame", False) else ""
    payload = json.dumps({
        "name": shape.name,
        "shape_id": shape.shape_id,
        "text": text,
        "bounds": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
    }, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(payload.encode()).hexdigest()


def _svg_realization(slide, semantics):
    diagram_ids = {
        obj["object_id"]
        for obj in semantics["semantic_objects"]
        if obj["role"] in {"diagram_node", "connector"}
    }
    parts = {object_id: [] for object_id in diagram_ids}
    prefix = f"oxq:{semantics['slide_id']}:"
    for shape in slide.shapes:
        name = shape.name or ""
        if not name.startswith(prefix):
            continue
        remainder = name[len(prefix):]
        if ":" not in remainder:
            continue
        object_id, part = remainder.split(":", 1)
        if object_id not in diagram_ids:
            continue
        parts[object_id].append({
            "ppt_object_id": str(shape.shape_id),
            "ppt_shape_name": name,
            "bounds_emu": {
                "x": int(shape.left), "y": int(shape.top),
                "w": int(shape.width), "h": int(shape.height),
            },
            "content_hash": _shape_signature(shape),
            "part": part,
        })
    missing = sorted(object_id for object_id, values in parts.items() if not values)
    if missing:
        raise RuntimeError(f"SVG lane missing semantic objects after trace stamping: {missing}")
    records = []
    for object_id, values in sorted(parts.items()):
        aggregate = "|".join(v["content_hash"] for v in values)
        role = object_index(semantics)[object_id]["role"]
        records.append({
            "semantic_object_id": object_id,
            "render_lane": "svg",
            "parts": values,
            "token_refs": object_index(semantics)[object_id].get("token_refs", []),
            "fidelity": "semantic_and_editable" if role == "diagram_node" else "editable_vector_structure",
            "fit_result": "fit",
            "fallback": None,
            "realization_rev": 1,
            "pptx_content_hash": hashlib.sha256(aggregate.encode()).hexdigest(),
        })
    return records


def _render_hybrid_problem(em, root):
    s = em.semantics
    em.text("title", _text(s, "title"), "title_band", (0, 0, 1, 0.52), role="title", part="title")
    em.text("subtitle", _text(s, "subtitle"), "title_band", (0, 0.58, 1, 0.32), role="subtitle", part="subtitle")
    em.metric_card("metric_wait", _text(s, "metric_wait"), _label(s, "metric_wait"), "problem_copy", (0.00, 0.04, 0.65, 0.36))
    em.text("pain_annotation", _text(s, "pain_annotation"), "problem_copy", (0.00, 0.50, 0.92, 0.38), role="annotation", part="annotation")
    hero = root / "experiment/queuezero/assets/problem_hero_structural_v1.png"
    if not hero.exists():
        raise RuntimeError("structural hero fixture missing; run scripts/create_queuezero_mock_assets.py")
    em.picture("hero_visual_slot", hero, "hero_frame", lane="image")
    em.text("source_note", _text(s, "source_note"), "footer_region", role="source_note", part="source")


def _decorate_how_it_works(em, root):
    s = em.semantics
    em.text("title", _text(s, "title"), "title_band", (0, 0, 1, 0.52), role="title", part="title")
    em.text("subtitle", _text(s, "subtitle"), "title_band", (0, 0.58, 1, 0.32), role="subtitle", part="subtitle")
    asset = next(a for a in s["assets"] if a["semantic_object_id"] == "screenshot_main")
    source = root / asset["current_instance"]["source"]
    em.picture("screenshot_main", source, "screenshot_region", lane="native")
    em.text("source_note", _text(s, "source_note"), "footer_region", role="source_note", part="source")


def _move_second_slide_to_front(prs):
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    if len(ids) != 3:
        raise RuntimeError(f"expected three slides before reorder, got {len(ids)}")
    sld_id_lst.remove(ids[1])
    sld_id_lst.insert(0, ids[1])


def build_hybrid(root, ppt_master_root, output_pptx, realization_dir, adapter_workspace=None):
    root = Path(root)
    deck_system = load_json(root / "experiment/queuezero/deck_system.stage3.json")
    semantics_dir = root / "experiment/queuezero/slide_semantics"
    problem = load_json(semantics_dir / "problem_hook.v1.json")
    how = load_json(semantics_dir / "how_it_works.v1.json")
    validation = load_json(semantics_dir / "validation_traction.v1.json")

    if adapter_workspace is None:
        temp = tempfile.TemporaryDirectory(prefix="queuezero-ppt-master-")
        workspace = Path(temp.name)
    else:
        temp = None
        workspace = Path(adapter_workspace)
        workspace.mkdir(parents=True, exist_ok=True)

    adapter = compile_how_it_works(root, ppt_master_root, workspace)
    prs = Presentation(adapter["pptx"])
    if len(prs.slides) != 1:
        raise RuntimeError("hybrid base must contain one PPT Master slide")
    blank = prs.slide_layouts[6]

    # Slide 1 currently holds the SVG-compiled How-It-Works diagram.
    how_slide = prs.slides[0]
    svg_records = _svg_realization(how_slide, how)
    how_em = Stage3Emitter(how, deck_system, how_slide, variant="hybrid")
    _decorate_how_it_works(how_em, root)
    how_native = how_em.realization()["objects"]
    how_realization = {
        "schema_version": "stage3-realization-v1",
        "slide_id": how["slide_id"],
        "variant": "hybrid",
        "semantic_hash": canonical_hash(how),
        "render_plan": "experiment/queuezero/render_plans/how_it_works.hybrid.json",
        "svg_adapter": {
            "ppt_master_commit": adapter["ppt_master_commit"],
            "svg_sha256": adapter["svg_sha256"],
            "trace": str(Path(adapter["trace"]).name),
        },
        "objects": svg_records + how_native,
    }

    problem_slide = prs.slides.add_slide(blank)
    problem_em = Stage3Emitter(problem, deck_system, problem_slide, variant="hybrid")
    _render_hybrid_problem(problem_em, root)
    problem_realization = problem_em.realization()
    problem_realization["semantic_hash"] = canonical_hash(problem)
    problem_realization["render_plan"] = "experiment/queuezero/render_plans/problem_hook.hybrid.json"
    problem_realization["visual_quality_status"] = "structural_fixture_only_not_scored"

    validation_slide = prs.slides.add_slide(blank)
    validation_em = Stage3Emitter(validation, deck_system, validation_slide, variant="hybrid")
    _render_validation(validation_em)
    validation_realization = validation_em.realization()
    validation_realization["semantic_hash"] = canonical_hash(validation)
    validation_realization["render_plan"] = "experiment/queuezero/render_plans/validation_traction.hybrid.json"

    # Current order is How-It-Works, Problem, Validation. Benchmark order is Problem, How, Validation.
    _move_second_slide_to_front(prs)

    output_pptx = Path(output_pptx)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)

    realization_dir = Path(realization_dir)
    realization_dir.mkdir(parents=True, exist_ok=True)
    realizations = [problem_realization, how_realization, validation_realization]
    for realization in realizations:
        path = realization_dir / f"{realization['slide_id']}.hybrid.json"
        path.write_text(json.dumps(realization, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    adapter_meta = realization_dir / "ppt_master_adapter.json"
    adapter_meta.write_text(json.dumps({
        "ppt_master_commit": adapter["ppt_master_commit"],
        "svg_sha256": adapter["svg_sha256"],
        "mapped": adapter["mapped"],
        "visual_quality_status": "problem hero uses structural fixture until real image-generation asset is supplied",
    }, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    if temp is not None:
        temp.cleanup()
    return output_pptx, realizations
