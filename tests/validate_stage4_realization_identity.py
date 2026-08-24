#!/usr/bin/env python3
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from tests.audit_stage4_deck import _shape_signature, validate_realization_identity


def _set_name(shape, name):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("shape naming metadata missing")


def _fixture():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    _set_name(shape, "oxq:problem-hook:title:title")
    shape.text = "Peak lunch queues waste time"
    record = {
        "slide_id": "problem-hook",
        "objects": [{
            "semantic_object_id": "title",
            "part": "title",
            "kind": "text",
            "ppt_shape_name": shape.name,
            "ppt_object_id": str(shape.shape_id),
            "bounds_emu": {
                "x": int(shape.left), "y": int(shape.top),
                "w": int(shape.width), "h": int(shape.height),
            },
            "content_hash": _shape_signature(shape),
            "render_lane": "native",
            "fidelity": "semantic_and_editable",
        }],
    }
    return slide, record


def _expect_rejection(slide, record, expected):
    try:
        validate_realization_identity(slide, record)
    except ValueError as exc:
        if expected not in str(exc):
            raise AssertionError(f"expected {expected!r}; got {exc!r}") from exc
        return
    raise AssertionError(f"invalid realization unexpectedly accepted: {expected}")


def main():
    slide, good = _fixture()
    validate_realization_identity(slide, good)

    changed = copy.deepcopy(good)
    changed["objects"][0]["content_hash"] = "0" * 64
    _expect_rejection(slide, changed, "content-hash drift")

    changed = copy.deepcopy(good)
    changed["objects"][0]["bounds_emu"]["w"] += 1
    _expect_rejection(slide, changed, "bounds drift")

    changed = copy.deepcopy(good)
    changed["objects"][0]["ppt_object_id"] = "999999"
    _expect_rejection(slide, changed, "object-id drift")

    extra = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(2), Inches(1))
    _set_name(extra, "oxq:problem-hook:subtitle:subtitle")
    extra.text = "Unrecorded semantic shape"
    _expect_rejection(slide, good, "semantic-object set drift")

    print("Stage 4 realization identity regression checks: PASS")


if __name__ == "__main__":
    main()
