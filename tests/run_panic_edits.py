#!/usr/bin/env python3
import argparse
import hashlib
import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "experiment/queuezero/assets"
ORANGE = RGBColor(0xF9, 0x73, 0x16)
BLUE_HEX = "2563EB"
PX30_EMU = 30 * 9525


def shape_name_set(prs):
    out = {}
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            name = shape.name or f"unnamed:{slide_idx}:{shape.shape_id}"
            payload = shape._element.xml.encode("utf-8")
            out[name] = hashlib.sha256(payload).hexdigest()
    return out


def find_exact(prs, name):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == name:
                return slide, shape
    raise KeyError(name)


def find_prefix(prs, prefix):
    return [
        (slide, shape)
        for slide in prs.slides
        for shape in slide.shapes
        if (shape.name or "").startswith(prefix)
    ]


def set_shape_name(shape, name):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("shape has no cNvPr")


def set_text_preserve(shape, text):
    if not shape.has_text_frame:
        raise RuntimeError(f"{shape.name} is not editable text")
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for run in p.runs[1:]:
            run.text = ""
    else:
        p.text = text
    for para in tf.paragraphs[1:]:
        para.text = ""


def remove_shape(shape):
    elem = shape._element
    elem.getparent().remove(elem)


def replace_picture(prs, name, source):
    slide, shape = find_exact(prs, name)
    x, y, w, h = shape.left, shape.top, shape.width, shape.height
    remove_shape(shape)
    new_shape = slide.shapes.add_picture(str(source), x, y, w, h)
    set_shape_name(new_shape, name)
    return new_shape


def sponsor_rect(prs):
    semantics = json.loads((ROOT / "experiment/queuezero/slide_semantics/problem_hook.v1.json").read_text(encoding="utf-8"))
    region = next(r for r in semantics["regions"] if r["region_id"] == "sponsor_insertion")["rect"]
    return (
        int(region["x"] * prs.slide_width),
        int(region["y"] * prs.slide_height),
        int(region["w"] * prs.slide_width),
        int(region["h"] * prs.slide_height),
    )


def add_sponsor(prs):
    slide = prs.slides[0]
    x, y, w, h = sponsor_rect(prs)
    shape = slide.shapes.add_picture(str(ASSETS / "sponsor_logo_v1.png"), x, y, w, h)
    set_shape_name(shape, "oxq:problem-hook:sponsor_logo:user_added")
    return shape


def recolor_blue(prs):
    touched = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            changed = False
            try:
                if str(shape.fill.fore_color.rgb) == BLUE_HEX:
                    shape.fill.fore_color.rgb = ORANGE
                    changed = True
            except Exception:
                pass
            try:
                if str(shape.line.color.rgb) == BLUE_HEX:
                    shape.line.color.rgb = ORANGE
                    changed = True
            except Exception:
                pass
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        try:
                            if str(run.font.color.rgb) == BLUE_HEX:
                                run.font.color.rgb = ORANGE
                                changed = True
                        except Exception:
                            pass
            if changed:
                touched.add(shape.name)
    return touched


def changed_names(before, after):
    names = set(before) | set(after)
    return {name for name in names if before.get(name) != after.get(name)}


def matches_expected(name, expected):
    return any(name == item or (item.endswith("*") and name.startswith(item[:-1])) for item in expected)


def run_one(original, out_path, edit_id, variant):
    prs = Presentation(original)
    before = shape_name_set(prs)
    expected = []
    outcome = "DIRECT_PPT_EDIT"
    score = 5
    notes = []

    if edit_id == "P1":
        name = "oxq:validation-traction:metric_weekly_intent:value"
        _, shape = find_exact(prs, name)
        set_text_preserve(shape, "81%")
        expected = [name]
    elif edit_id == "P2":
        name = "oxq:problem-hook:title:title"
        _, shape = find_exact(prs, name)
        set_text_preserve(shape, "Peak lunch queues create a daily time tax students cannot predict")
        expected = [name]
    elif edit_id == "P3":
        name = "oxq:how-it-works:screenshot_main:picture"
        replace_picture(prs, name, ASSETS / "product_ui_v2.png")
        expected = [name]
    elif edit_id == "P4":
        prefix = "oxq:validation-traction:metric_weekly_intent:"
        matches = find_prefix(prs, prefix)
        if not matches:
            raise RuntimeError("weekly-intent KPI not found")
        for _, shape in matches:
            shape.left += PX30_EMU
        expected = [prefix + "*"]
    elif edit_id == "P5":
        name = "oxq:problem-hook:title:title"
        _, shape = find_exact(prs, name)
        sizes = []
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    sizes.append(run.font.size)
                    run.font.size = int(run.font.size * 0.85)
        if not sizes:
            raise RuntimeError("title has no explicit font size")
        expected = [name]
    elif edit_id == "P6":
        touched = recolor_blue(prs)
        expected = list(touched)
        if variant == "hybrid":
            outcome, score = "LOCAL_ASSET_REGEN", 4
            notes.append("bounded Problem hero declares accent.primary dependency and would require local asset regeneration for full visual consistency")
    elif edit_id == "P7":
        prefix = "oxq:validation-traction:metric_cafeterias:"
        matches = find_prefix(prs, prefix)
        if not matches:
            raise RuntimeError("secondary cafeteria metric not found")
        for _, shape in matches:
            remove_shape(shape)
        expected = [prefix + "*"]
        notes.append("deletion intentionally leaves layout hole; reflow quality is scored separately")
    elif edit_id == "P8":
        add_sponsor(prs)
        expected = ["oxq:problem-hook:sponsor_logo:user_added"]
    elif edit_id == "P9":
        name = "oxq:how-it-works:node_wait_predictor:label"
        _, shape = find_exact(prs, name)
        set_text_preserve(shape, "Predicted wait time")
        expected = [name]
    elif edit_id == "P10":
        if variant != "hybrid":
            return {
                "edit_id": edit_id,
                "variant": variant,
                "applicable": False,
                "outcome": "NOT_APPLICABLE",
                "score": None,
                "notes": ["native-vector baseline has no bounded generated hero asset"],
            }
        name = "oxq:problem-hook:hero_visual_slot:picture"
        replace_picture(prs, name, ASSETS / "problem_hero_structural_v2.png")
        expected = [name]
        outcome, score = "LOCAL_ASSET_REGEN", 4
        notes.append("structural fixture proves blast radius only; this is not visual-quality evidence")
    else:
        raise ValueError(edit_id)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    reopened = Presentation(out_path)
    after = shape_name_set(reopened)
    changed = changed_names(before, after)
    unexpected = sorted(name for name in changed if not matches_expected(name, expected))
    expected_changed = sorted(name for name in changed if matches_expected(name, expected))
    if unexpected:
        score = min(score, 3)
        notes.append(f"unexpected changed objects: {unexpected}")
    collateral = len(unexpected) / max(1, len(before))
    return {
        "edit_id": edit_id,
        "variant": variant,
        "applicable": True,
        "outcome": outcome,
        "score": score,
        "pre_edit_addressable_objects": len(before),
        "changed_objects": sorted(changed),
        "expected_changed_objects": expected_changed,
        "unexpected_changed_objects": unexpected,
        "collateral_damage_ratio": collateral,
        "ai_call_count": 0 if edit_id != "P10" and not (edit_id == "P6" and variant == "hybrid") else None,
        "notes": notes,
    }


def p8_then_p1(original, variant, out_dir):
    stage1 = out_dir / "P8_then_P1_stage1.pptx"
    prs = Presentation(original)
    add_sponsor(prs)
    prs.save(stage1)
    before = shape_name_set(Presentation(stage1))
    prs = Presentation(stage1)
    _, metric = find_exact(prs, "oxq:validation-traction:metric_weekly_intent:value")
    set_text_preserve(metric, "81%")
    final = out_dir / "P8_then_P1_final.pptx"
    prs.save(final)
    reopened = Presentation(final)
    after = shape_name_set(reopened)
    changed = changed_names(before, after)
    sponsor_survives = "oxq:problem-hook:sponsor_logo:user_added" in after
    expected = {"oxq:validation-traction:metric_weekly_intent:value"}
    return {
        "sequence": "P8_then_P1",
        "variant": variant,
        "sponsor_survives_reopen_and_later_edit": sponsor_survives,
        "changed_during_P1": sorted(changed),
        "unexpected_changed_objects": sorted(changed - expected),
        "passed": sponsor_survives and changed == expected,
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--variant", choices=["native_vector", "hybrid"], required=True)
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--out-dir", required=True)
    args = parser.parse_args()
    original = ROOT / args.pptx
    out_dir = ROOT / args.out_dir
    if not original.exists():
        raise SystemExit(f"missing input deck: {original}")
    out_dir.mkdir(parents=True, exist_ok=True)
    results = []
    for edit_id in [f"P{i}" for i in range(1, 11)]:
        results.append(run_one(original, out_dir / f"{edit_id}.pptx", edit_id, args.variant))
    sequence = p8_then_p1(original, args.variant, out_dir)
    summary = {
        "variant": args.variant,
        "results": results,
        "sequence_interaction": sequence,
        "editability_score_sum": sum(r["score"] for r in results if r.get("score") is not None),
        "editability_score_max": 5 * sum(1 for r in results if r.get("score") is not None),
    }
    summary["normalized_editability"] = summary["editability_score_sum"] / max(1, summary["editability_score_max"])
    path = out_dir / "summary.json"
    path.write_text(json.dumps(summary, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    if any(r.get("unexpected_changed_objects") for r in results if r.get("applicable")):
        raise SystemExit(f"panic edits found collateral changes; see {path}")
    if not sequence["passed"]:
        raise SystemExit(f"P8→P1 interaction failed; see {path}")
    print(json.dumps(summary, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
