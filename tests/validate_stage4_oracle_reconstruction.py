#!/usr/bin/env python3
"""Fail-closed checks for the Stage 4 oracle reconstruction lane."""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "dist/stage4/oracle_reconstruction_v1"
REALIZATIONS = OUT / "realizations.json"
PPTX = OUT / "queuezero_stage4_oracle_reconstruction_v1.pptx"

if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from src.ir.runtime import load_json, canonical_hash


def fail(message):
    raise RuntimeError(message)


def semantic_expectations(semantic_path: Path):
    sem = load_json(semantic_path)
    editable = {
        item["object_id"]
        for item in sem.get("semantic_objects", [])
        if item.get("editability_priority") == "must_remain_editable"
    }
    replaceable = {
        item["object_id"]
        for item in sem.get("semantic_objects", [])
        if item.get("editability_priority") == "replaceable_asset"
    }
    return sem, editable, replaceable


def main():
    if not REALIZATIONS.is_file() or not PPTX.is_file():
        fail("oracle reconstruction outputs are missing; run build script first")
    payload = json.loads(REALIZATIONS.read_text(encoding="utf-8"))
    if payload.get("full_slide_oracle_raster_used") is not False:
        fail("full-slide oracle raster shortcut is forbidden")
    if payload.get("slide_count") != 3:
        fail("expected exactly three reconstruction slides")

    prs = Presentation(PPTX)
    if len(prs.slides) != 3:
        fail("PPTX slide count mismatch")

    for slide_idx, realization in enumerate(payload.get("slides", [])):
        semantic_path = ROOT / realization["semantic_file"] if "semantic_file" in realization else None
        if semantic_path is None:
            # realization carries the hash, while the derived IR record carries the semantic path indirectly
            derived = payload["derived_visual_ir"][slide_idx]
            ir = load_json(ROOT / derived["path"])
            semantic_path = ROOT / ir["semantic_file"]
        sem, editable, replaceable = semantic_expectations(semantic_path)
        if canonical_hash(sem) != realization.get("semantic_hash"):
            fail(f"semantic hash mismatch on {realization['slide_id']}")

        records = realization.get("objects", [])
        realized_ids = {item.get("semantic_object_id") for item in records}
        missing = sorted(editable - realized_ids)
        if missing:
            fail(f"editable semantic objects missing on {realization['slide_id']}: {missing}")

        for item in records:
            kind = item.get("kind")
            oid = item.get("semantic_object_id")
            if kind == "picture":
                if oid not in replaceable:
                    fail(f"picture lane used for non-replaceable semantic object {oid}")
                bounds = item["bounds_emu"]
                area_fraction = (bounds["w"] * bounds["h"]) / (prs.slide_width * prs.slide_height)
                if area_fraction >= 0.50:
                    fail(f"picture {oid} occupies {area_fraction:.3f} of slide; full-slide shortcut risk")
                asset_path = item.get("asset_path", "")
                if asset_path.endswith(("problem-hook.jpg", "how-it-works.jpg", "validation-traction.jpg")):
                    fail(f"raw full-slide oracle image used directly: {asset_path}")
            elif kind in {"text", "line", "connector"}:
                if item.get("fidelity") != "semantic_and_editable":
                    fail(f"native object lost editable fidelity: {oid}/{kind}")

        # Verify every must-remain-editable object has at least one named PPT shape.
        names = {shape.name for shape in prs.slides[slide_idx].shapes}
        for oid in editable:
            prefix = f"oxq:{realization['slide_id']}:{oid}:"
            if not any(name.startswith(prefix) for name in names):
                fail(f"PPTX missing named editable object {oid} on {realization['slide_id']}")

    bindings = json.loads((OUT / "asset_bindings.json").read_text(encoding="utf-8"))
    for asset in bindings.get("derived_assets", []):
        if asset.get("operation") != "deterministic_text_free_crop":
            fail("oracle-derived raster must use declared deterministic text-free crop operation")
        if asset.get("semantic_authority") is not False:
            fail("derived raster must never claim semantic authority")

    print("Stage 4 oracle reconstruction structural/editability gate: PASS")


if __name__ == "__main__":
    main()
