#!/usr/bin/env python3
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "dist/queuezero_hybrid.pptx"
REAL = ROOT / "dist/realizations/hybrid"


def fail(message):
    raise SystemExit(f"HYBRID AUDIT FAILED: {message}")


def named(slide):
    return {shape.name: shape for shape in slide.shapes}


def is_full_slide_picture(shape, prs):
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    return shape.width >= prs.slide_width * 0.95 and shape.height >= prs.slide_height * 0.95


def main():
    if not PPTX.exists():
        fail(f"missing {PPTX}")
    prs = Presentation(PPTX)
    if len(prs.slides) != 3:
        fail(f"expected 3 slides, got {len(prs.slides)}")

    for idx, slide in enumerate(prs.slides, 1):
        offenders = [shape.name for shape in slide.shapes if is_full_slide_picture(shape, prs)]
        if offenders:
            fail(f"slide {idx} contains full-slide picture(s): {offenders}")

    problem, how, validation = prs.slides[0], prs.slides[1], prs.slides[2]
    problem_names = named(problem)
    how_names = named(how)

    hero_name = "oxq:problem-hook:hero_visual_slot:picture"
    if hero_name not in problem_names or problem_names[hero_name].shape_type != MSO_SHAPE_TYPE.PICTURE:
        fail("Problem bounded hero slot is not an independently addressable picture")

    screenshot_name = "oxq:how-it-works:screenshot_main:picture"
    if screenshot_name not in how_names or how_names[screenshot_name].shape_type != MSO_SHAPE_TYPE.PICTURE:
        fail("How-It-Works screenshot is not an independently addressable picture")

    diagram_objects = [
        "node_camera", "node_queue_estimator", "node_wait_predictor", "node_decision",
        "connector_camera_queue", "connector_queue_prediction", "connector_prediction_decision",
    ]
    for object_id in diagram_objects:
        prefix = f"oxq:how-it-works:{object_id}:"
        matches = [shape for name, shape in how_names.items() if name.startswith(prefix)]
        if not matches:
            fail(f"missing SVG-lane semantic object {object_id}")
        if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in matches):
            fail(f"SVG-lane semantic object {object_id} rasterized to picture")

    label_text = {
        name: shape.text
        for name, shape in how_names.items()
        if name.startswith("oxq:how-it-works:node_") and getattr(shape, "has_text_frame", False)
    }
    expected_fragments = ["Camera feed", "Queue length", "Wait prediction", "Go now / choose", "another cafeteria"]
    combined = "\n".join(label_text.values())
    missing_text = [text for text in expected_fragments if text not in combined]
    if missing_text:
        fail(f"SVG lane did not preserve editable native text: {missing_text}")

    for slide_id in ["problem-hook", "how-it-works", "validation-traction"]:
        path = REAL / f"{slide_id}.hybrid.json"
        if not path.exists():
            fail(f"missing realization {path.name}")
        realization = json.loads(path.read_text(encoding="utf-8"))
        if realization.get("variant") != "hybrid":
            fail(f"wrong variant in {path.name}")

    how_real = json.loads((REAL / "how-it-works.hybrid.json").read_text(encoding="utf-8"))
    by_id = {obj["semantic_object_id"]: obj for obj in how_real["objects"]}
    for object_id in diagram_objects:
        if by_id.get(object_id, {}).get("render_lane") != "svg":
            fail(f"{object_id} realization is not recorded as SVG lane")
    if by_id.get("screenshot_main", {}).get("render_lane") != "native":
        fail("screenshot realization should be native picture slot")

    problem_real = json.loads((REAL / "problem-hook.hybrid.json").read_text(encoding="utf-8"))
    problem_by_id = {obj["semantic_object_id"]: obj for obj in problem_real["objects"]}
    if problem_by_id.get("hero_visual_slot", {}).get("render_lane") != "image":
        fail("Problem hero is not recorded as bounded image lane")

    print("Hybrid audit passed")
    print(f"slides={len(prs.slides)}; how_shapes={len(how.shapes)}; editable_svg_labels={len(label_text)}")


if __name__ == "__main__":
    main()
