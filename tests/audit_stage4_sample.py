#!/usr/bin/env python3
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "dist/stage4/validation_traction_sample_v0.pptx"
SOLUTION = ROOT / "dist/stage4/validation_traction.layout_solution.v0.json"
REALIZATION = ROOT / "dist/stage4/validation_traction.realization.v0.json"

EXPECTED_TEXT = {
    "oxq:validation-traction:title:title": "Early testing clears the bar for a semester pilot",
    "oxq:validation-traction:subtitle:subtitle": "Technical accuracy and user intent are promising; retention and paid adoption remain unproven.",
    "oxq:validation-traction:metric_weekly_intent:value": "76%",
    "oxq:validation-traction:metric_weekly_intent:label": "would use weekly",
    "oxq:validation-traction:metric_mae:value": "3.8 min",
    "oxq:validation-traction:metric_mae:label": "prediction MAE",
    "oxq:validation-traction:metric_students:value": "42",
    "oxq:validation-traction:metric_students:label": "students tested",
    "oxq:validation-traction:metric_cafeterias:value": "3",
    "oxq:validation-traction:metric_cafeterias:label": "cafeterias",
    "oxq:validation-traction:pilot_gate:label": "Next proof: one-university semester pilot",
    "oxq:validation-traction:source_note:source": "Source: QueueZero controlled benchmark brief",
}

FORBIDDEN_VISIBLE = {
    "metric_weekly_intent", "metric_mae", "validation_stack", "validation_evidence_stack",
    "stage4", "visual_ir", "semantic_object_id", "oxq:"
}


def main():
    for path in (PPTX, SOLUTION, REALIZATION):
        if not path.exists():
            raise SystemExit(f"missing Stage 4 sample artifact: {path}")

    prs = Presentation(PPTX)
    if len(prs.slides) != 1:
        raise SystemExit(f"expected one sample slide, got {len(prs.slides)}")
    slide = prs.slides[0]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if pictures:
        raise SystemExit(f"Validation sample must use zero pictures; found {len(pictures)}")

    by_name = {shape.name: shape for shape in slide.shapes}
    missing = sorted(set(EXPECTED_TEXT) - set(by_name))
    if missing:
        raise SystemExit(f"missing semantic shapes: {missing}")

    for name, expected in EXPECTED_TEXT.items():
        shape = by_name[name]
        if not shape.has_text_frame:
            raise SystemExit(f"semantic shape is not editable PowerPoint text: {name}")
        actual = shape.text.strip()
        if actual != expected:
            raise SystemExit(f"semantic text mismatch {name}: {actual!r} != {expected!r}")

    visible_text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    lower_visible = visible_text.lower()
    leaked = sorted(token for token in FORBIDDEN_VISIBLE if token.lower() in lower_visible)
    if leaked:
        raise SystemExit(f"internal IR/schema labels leaked into audience-visible text: {leaked}")

    hero = by_name["oxq:validation-traction:metric_weekly_intent:value"]
    hero_runs = [run for p in hero.text_frame.paragraphs for run in p.runs]
    if not hero_runs or hero_runs[0].font.size is None or hero_runs[0].font.size.pt < 70:
        raise SystemExit("76% protagonist is not rendered at hero scale")

    title = by_name["oxq:validation-traction:title:title"]
    subtitle = by_name["oxq:validation-traction:subtitle:subtitle"]
    if int(title.top + title.height) > int(subtitle.top):
        raise SystemExit("title/subtitle geometry collides")

    semantic_rectangles = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.name.startswith("oxq:validation-traction:metric_")
    ]
    if semantic_rectangles:
        raise SystemExit("metric evidence regressed into semantic metric-card rectangles")

    solution = json.loads(SOLUTION.read_text(encoding="utf-8"))
    realization = json.loads(REALIZATION.read_text(encoding="utf-8"))
    if solution["semantic_hash"] != realization["semantic_hash"]:
        raise SystemExit("semantic hash drift between compiler solution and realization")
    if solution["visual_ir_hash"] != realization["visual_ir_hash"]:
        raise SystemExit("Visual IR hash drift between compiler solution and realization")
    if len(realization.get("objects", [])) != len(EXPECTED_TEXT):
        raise SystemExit("realization does not map every editable semantic text part")

    print("Stage 4 Validation sample structural audit passed")
    print(f"  slides: 1")
    print(f"  editable semantic text parts: {len(EXPECTED_TEXT)}")
    print(f"  pictures: {len(pictures)}")
    print(f"  semantic metric-card shapes: {len(semantic_rectangles)}")


if __name__ == "__main__":
    main()
