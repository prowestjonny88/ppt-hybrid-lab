#!/usr/bin/env python3
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "dist" / "queuezero_native_vector.pptx"
REALIZATION_DIR = ROOT / "dist" / "realizations" / "native_vector"
EXPECTED_SLIDES = ["problem-hook", "how-it-works", "validation-traction"]


def main():
    if not PPTX.exists():
        print(f"Missing {PPTX}", file=sys.stderr)
        return 2
    prs = Presentation(PPTX)
    errors = []
    if len(prs.slides) != 3:
        errors.append(f"expected 3 slides, got {len(prs.slides)}")

    total_text = 0
    total_pictures = 0
    full_slide_pictures = 0
    all_names = []

    for idx, slide in enumerate(prs.slides):
        names = []
        for shape in slide.shapes:
            names.append(shape.name)
            all_names.append(shape.name)
            if not shape.name.startswith("oxq:"):
                errors.append(f"slide {idx+1}: untracked shape identity {shape.name!r}")
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                total_text += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_pictures += 1
                if shape.width >= prs.slide_width * 0.95 and shape.height >= prs.slide_height * 0.95:
                    full_slide_pictures += 1
        if len(names) != len(set(names)):
            errors.append(f"slide {idx+1}: duplicate deterministic shape names")

    if full_slide_pictures:
        errors.append(f"native-vector deck contains {full_slide_pictures} full-slide picture(s)")
    if total_text < 15:
        errors.append(f"expected substantial editable text; found only {total_text} text shapes")
    if total_pictures != 1:
        errors.append(f"expected exactly one replaceable screenshot picture; found {total_pictures}")

    for slide_id in EXPECTED_SLIDES:
        path = REALIZATION_DIR / f"{slide_id}.native_vector.json"
        if not path.exists():
            errors.append(f"missing realization {path.name}")
            continue
        data = json.loads(path.read_text(encoding="utf-8"))
        if data.get("slide_id") != slide_id or data.get("variant") != "native_vector":
            errors.append(f"invalid realization header for {slide_id}")
        for record in data.get("objects", []):
            if not record.get("parts"):
                errors.append(f"{slide_id}:{record.get('semantic_object_id')} has no PPT parts")
            for part in record.get("parts", []):
                if part.get("ppt_shape_name") not in all_names:
                    errors.append(f"realization references missing shape {part.get('ppt_shape_name')}")

    if errors:
        print("Native output audit FAILED", file=sys.stderr)
        for error in errors:
            print(f"  - {error}", file=sys.stderr)
        return 1

    print(json.dumps({
        "status": "PASS",
        "slides": len(prs.slides),
        "editable_text_shapes": total_text,
        "pictures": total_pictures,
        "full_slide_pictures": full_slide_pictures,
        "tracked_shapes": len(all_names)
    }, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
