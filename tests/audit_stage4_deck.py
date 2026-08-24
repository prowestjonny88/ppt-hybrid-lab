#!/usr/bin/env python3
import hashlib
import json
import math
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from src.ir.runtime import canonical_hash, load_json
from src.visual_ir.style_runtime import resolve_style_profile

PPTX = ROOT / "dist/stage4/deck/queuezero_stage4_v0.pptx"
REAL = ROOT / "dist/stage4/deck/realizations.json"
LAYOUT = ROOT / "dist/stage4/deck/layout_solutions.json"
STYLE_PATH = ROOT / "experiment/queuezero/style_profiles/queuezero_hackathon_v0.json"
EXPECTED_SLIDES = ["problem-hook", "how-it-works", "validation-traction"]
VISUAL_IR_BY_SLIDE = {
    "problem-hook": ROOT / "experiment/queuezero/visual_ir/problem_hook.stage4.v0.json",
    "how-it-works": ROOT / "experiment/queuezero/visual_ir/how_it_works.stage4.v0.json",
    "validation-traction": ROOT / "experiment/queuezero/visual_ir/validation_traction.stage4.v0.json",
}
EMU_PER_INCH = 914400


def _max_font_pt(shape):
    sizes = []
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return max(sizes) if sizes else None


def _shape_signature(shape):
    text = shape.text if getattr(shape, "has_text_frame", False) else ""
    payload = json.dumps(
        {
            "shape_id": shape.shape_id,
            "name": shape.name,
            "text": text,
            "bounds": [
                int(shape.left),
                int(shape.top),
                int(shape.width),
                int(shape.height),
            ],
        },
        sort_keys=True,
        separators=(",", ":"),
    )
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _estimated_required_height_in(shape):
    """Conservative wrap estimate for high-risk semantic text.

    This is deliberately a proxy, not a replacement for pixel review. It catches
    the failure class where the PowerPoint textbox bounds are valid but a renderer
    wraps glyphs beyond the intended semantic zone and into adjacent content.
    """
    text = " ".join((shape.text or "").split())
    font_pt = _max_font_pt(shape)
    if not text or not font_pt:
        return 0.0
    width_in = shape.width / EMU_PER_INCH
    chars_per_line = max(1, int((width_in * 72.0) / (font_pt * 0.62)))
    line_count = max(1, math.ceil(len(text) / chars_per_line))
    return line_count * font_pt * 1.20 / 72.0


def _check_text_capacity(slide, slide_id):
    risky = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        font_pt = _max_font_pt(shape)
        is_title_band = (
            shape.name.endswith(":title:title")
            or shape.name.endswith(":subtitle:subtitle")
        )
        is_large_evidence = font_pt is not None and font_pt >= 24
        if not (is_title_band or is_large_evidence):
            continue
        required = _estimated_required_height_in(shape)
        available = shape.height / EMU_PER_INCH
        if required > available * 0.96:
            risky.append(
                f"{shape.name} requires~{required:.2f}in but box is {available:.2f}in"
            )
    if risky:
        raise ValueError(
            f"semantic text overflow risk on {slide_id}: " + "; ".join(risky)
        )


def validate_realization_identity(slide, realization):
    """Fail closed if realization metadata drifts from the actual PPTX objects."""
    slide_id = realization["slide_id"]
    by_name = {shape.name: shape for shape in slide.shapes}
    expected_semantic_names = set()

    if len(realization.get("objects", [])) != len({o["ppt_shape_name"] for o in realization.get("objects", [])}):
        raise ValueError(f"duplicate realization object names on {slide_id}")

    for obj in realization["objects"]:
        name = obj["ppt_shape_name"]
        expected_semantic_names.add(name)
        shape = by_name.get(name)
        if shape is None:
            raise ValueError(f"realization object missing from PPTX on {slide_id}: {name}")

        if str(shape.shape_id) != str(obj["ppt_object_id"]):
            raise ValueError(
                f"realization object-id drift on {slide_id}: {name} "
                f"recorded={obj['ppt_object_id']} actual={shape.shape_id}"
            )

        bounds = obj["bounds_emu"]
        actual_bounds = {
            "x": int(shape.left),
            "y": int(shape.top),
            "w": int(shape.width),
            "h": int(shape.height),
        }
        if bounds != actual_bounds:
            raise ValueError(
                f"realization bounds drift on {slide_id}: {name} "
                f"recorded={bounds} actual={actual_bounds}"
            )

        actual_hash = _shape_signature(shape)
        if obj["content_hash"] != actual_hash:
            raise ValueError(
                f"realization content-hash drift on {slide_id}: {name}"
            )

        if obj["kind"] == "picture" and not obj.get("asset_present"):
            raise ValueError(f"required visual asset missing: {obj.get('asset_path')}")

    actual_semantic_names = {
        shape.name
        for shape in slide.shapes
        if shape.name.startswith(f"oxq:{slide_id}:") and ":decor:" not in shape.name
    }
    if actual_semantic_names != expected_semantic_names:
        missing = sorted(expected_semantic_names - actual_semantic_names)
        unrecorded = sorted(actual_semantic_names - expected_semantic_names)
        raise ValueError(
            f"realization semantic-object set drift on {slide_id}: "
            f"missing={missing} unrecorded={unrecorded}"
        )


def validate_source_provenance(realization, layout_solution):
    """Prove the rendered realization belongs to current source IR and style inputs."""
    slide_id = realization["slide_id"]
    if layout_solution.get("slide_id") != slide_id:
        raise ValueError(
            f"layout/realization slide-id drift: layout={layout_solution.get('slide_id')} realization={slide_id}"
        )

    passthrough = [
        "semantic_hash",
        "visual_ir_hash",
        "design_language_id",
        "design_language_hash",
        "style_profile_id",
        "style_profile_hash",
        "resolved_style_hash",
        "archetype_id",
        "variant",
        "solver",
        "compiler",
    ]
    for key in passthrough:
        if realization.get(key) != layout_solution.get(key):
            raise ValueError(
                f"layout/realization provenance drift on {slide_id}: {key} "
                f"layout={layout_solution.get(key)!r} realization={realization.get(key)!r}"
            )

    visual_ir_path = VISUAL_IR_BY_SLIDE.get(slide_id)
    if visual_ir_path is None or not visual_ir_path.is_file():
        raise ValueError(f"no canonical Visual IR source registered for {slide_id}")
    visual_ir = load_json(visual_ir_path)
    semantic_path = ROOT / visual_ir["semantic_file"]
    semantics = load_json(semantic_path)
    resolved_style, profile, language = resolve_style_profile(ROOT, STYLE_PATH)

    expected = {
        "semantic_hash": canonical_hash(semantics),
        "visual_ir_hash": canonical_hash(visual_ir),
        "design_language_id": language["design_language_id"],
        "design_language_hash": canonical_hash(language),
        "style_profile_id": profile["profile_id"],
        "style_profile_hash": canonical_hash(profile),
        "resolved_style_hash": canonical_hash(resolved_style),
        "archetype_id": visual_ir["composition"]["archetype_id"],
        "variant": visual_ir["composition"]["variant"],
    }
    for key, expected_value in expected.items():
        if layout_solution.get(key) != expected_value:
            raise ValueError(
                f"source/layout provenance drift on {slide_id}: {key} "
                f"expected={expected_value!r} actual={layout_solution.get(key)!r}"
            )


def audit_deck(pptx_path=PPTX, realization_path=REAL, layout_path=LAYOUT):
    pptx_path = Path(pptx_path)
    realization_path = Path(realization_path)
    layout_path = Path(layout_path)
    if not pptx_path.is_file() or not realization_path.is_file() or not layout_path.is_file():
        raise ValueError("Stage 4 deck outputs missing")

    prs = Presentation(pptx_path)
    data = json.loads(realization_path.read_text(encoding="utf-8"))
    solutions = json.loads(layout_path.read_text(encoding="utf-8"))
    if len(prs.slides) != 3 or data.get("slide_count") != 3 or len(solutions) != 3:
        raise ValueError("Stage 4 deck must contain exactly three slides")
    ids = [s["slide_id"] for s in data["slides"]]
    if ids != EXPECTED_SLIDES:
        raise ValueError(f"unexpected slide order: {ids}")
    solution_ids = [s["slide_id"] for s in solutions]
    if solution_ids != EXPECTED_SLIDES:
        raise ValueError(f"unexpected layout solution order: {solution_ids}")

    all_names = []
    for slide, realization, solution in zip(prs.slides, data["slides"], solutions):
        names = [shape.name for shape in slide.shapes]
        if len(names) != len(set(names)):
            raise ValueError(f"duplicate PowerPoint shape names on {realization['slide_id']}")
        if not any(name.startswith(f"oxq:{realization['slide_id']}:") for name in names):
            raise ValueError(f"missing semantic identity on {realization['slide_id']}")
        all_names.extend(names)
        validate_source_provenance(realization, solution)
        validate_realization_identity(slide, realization)
        _check_text_capacity(slide, realization["slide_id"])

    if len(all_names) < 20:
        raise ValueError("deck unexpectedly sparse; structural realization likely incomplete")
    return True


def main():
    try:
        audit_deck()
    except ValueError as exc:
        raise SystemExit(str(exc)) from exc
    print("Stage 4 deck structural + provenance + realization identity audit: PASS")


if __name__ == "__main__":
    main()
