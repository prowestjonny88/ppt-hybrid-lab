#!/usr/bin/env python3
import hashlib
import json
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

TARGET_NAME = "oxq:validation-traction:metric_weekly_intent:value"


def set_shape_name(shape, name):
    # python-pptx exposes cNvPr through the shape XML; write the deterministic identity there.
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("cNvPr not found")


def add_metric(slide, name, value, x):
    box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.1), Inches(1.1))
    set_shape_name(box, name)
    p = box.text_frame.paragraphs[0]
    p.text = value
    p.runs[0].font.size = Pt(28)
    return box


def signature(shape):
    text = shape.text if getattr(shape, "has_text_frame", False) else None
    payload = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
        "text": text,
    }
    raw = json.dumps(payload, sort_keys=True, separators=(",", ":")).encode()
    payload["content_hash"] = hashlib.sha256(raw).hexdigest()
    return payload


def inspect(path):
    prs = Presentation(path)
    return {shape.name: signature(shape) for shape in prs.slides[0].shapes}


def build_baseline(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_metric(slide, TARGET_NAME, "76%", 1.0)
    add_metric(slide, "oxq:validation-traction:metric_mae:value", "3.8 min", 4.0)
    add_metric(slide, "oxq:validation-traction:metric_students:value", "42", 7.0)
    prs.save(path)


def patch_target(src, dst):
    prs = Presentation(src)
    matches = [shape for shape in prs.slides[0].shapes if shape.name == TARGET_NAME]
    if len(matches) != 1:
        raise AssertionError(f"expected exactly one target shape, found {len(matches)}")
    shape = matches[0]
    before_geometry = (shape.left, shape.top, shape.width, shape.height)
    shape.text_frame.paragraphs[0].text = "81%"
    after_geometry = (shape.left, shape.top, shape.width, shape.height)
    if before_geometry != after_geometry:
        raise AssertionError("target geometry changed during text patch")
    prs.save(dst)


def main():
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        baseline = td / "baseline.pptx"
        patched = td / "patched.pptx"
        build_baseline(baseline)
        before = inspect(baseline)
        patch_target(baseline, patched)
        after = inspect(patched)

        if set(before) != set(after):
            raise AssertionError("shape identity set changed after local patch")
        if before[TARGET_NAME]["text"] != "76%" or after[TARGET_NAME]["text"] != "81%":
            raise AssertionError("target value did not patch 76% -> 81%")

        target_before = dict(before[TARGET_NAME])
        target_after = dict(after[TARGET_NAME])
        target_before.pop("text")
        target_after.pop("text")
        target_before.pop("content_hash")
        target_after.pop("content_hash")
        if target_before != target_after:
            raise AssertionError("target changed outside its text content")

        changed = []
        for name in sorted(before):
            if before[name] != after[name]:
                changed.append(name)
            if name != TARGET_NAME and before[name] != after[name]:
                raise AssertionError(f"collateral mutation detected in {name}")

        if changed != [TARGET_NAME]:
            raise AssertionError(f"expected only target to change, got {changed}")

        result = {
            "probe": "V2_native_single_object_patch",
            "status": "PASS",
            "identity_convention": "oxq:{slide_id}:{object_id}[:{part}]",
            "requested_edit": "76% -> 81%",
            "changed_semantic_objects": ["metric_weekly_intent"],
            "collateral_semantic_object_changes": 0,
            "target_geometry_preserved": True,
            "shape_identity_preserved": True,
        }
        print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
