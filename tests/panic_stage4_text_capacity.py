#!/usr/bin/env python3
"""Fail-closed regression coverage for the Stage 4 semantic text-capacity audit."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from tests.audit_stage4_deck import _check_text_capacity


def _set_name(shape, name: str):
    for elem in shape._element.iter():
        if elem.tag.endswith("}cNvPr"):
            elem.set("name", name)
            return
    raise RuntimeError("unable to name test shape")


def _add_textbox(slide, name: str, text: str, width_in: float, height_in: float, font_pt: float):
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(width_in), Inches(height_in))
    _set_name(shape, name)
    tf = shape.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    return shape


def _expect_fail(slide, label: str):
    try:
        _check_text_capacity(slide, label)
    except SystemExit as exc:
        if "overflow risk" not in str(exc):
            raise SystemExit(f"panic produced wrong failure: {exc}")
        return
    raise SystemExit("text-capacity panic failed open")


def main():
    # Title-band regression: long title forced into a tiny box must be rejected.
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(
        slide,
        "oxq:panic:title:title",
        "This intentionally overlong title must wrap far beyond its declared semantic title band",
        1.4,
        0.28,
        30,
    )
    _expect_fail(slide, "panic-title")

    # Large evidence regression: even when it is not a title, oversized evidence
    # must not silently wrap into adjacent content.
    prs2 = Presentation()
    slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
    _add_textbox(
        slide2,
        "oxq:panic:metric:value",
        "12345 67890",
        0.7,
        0.30,
        42,
    )
    _expect_fail(slide2, "panic-evidence")

    # Control case: a normal short title in a roomy box should remain accepted.
    prs3 = Presentation()
    slide3 = prs3.slides.add_slide(prs3.slide_layouts[6])
    _add_textbox(slide3, "oxq:control:title:title", "Short title", 5.0, 0.8, 26)
    _check_text_capacity(slide3, "control")

    print("Stage 4 text-capacity panic suite: PASS")


if __name__ == "__main__":
    main()
